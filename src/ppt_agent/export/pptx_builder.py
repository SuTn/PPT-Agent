from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

SLIDE_WIDTH = Emu(12192000)   # 13.333 inches
SLIDE_HEIGHT = Emu(6858000)   # 7.5 inches


def build_pptx(png_paths: list[str], output_path: Path) -> list[str]:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]
    failed: list[str] = []

    for png_path in png_paths:
        slide = prs.slides.add_slide(blank_layout)
        try:
            slide.shapes.add_picture(png_path, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
        except Exception:
            failed.append(png_path)

    if not failed:
        prs.save(str(output_path))
    else:
        prs.save(str(output_path))

    return failed
